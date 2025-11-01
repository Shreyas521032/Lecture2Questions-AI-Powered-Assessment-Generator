import streamlit as st
from io import BytesIO
import google.generativeai as genai
import pandas as pd
from datetime import datetime
import plotly.express as px

# Set page configuration
st.set_page_config(
    page_title="Lecture2Exam",
    page_icon="🧠",
    layout="wide",
    initial_sidebar_state="collapsed"
)

st.markdown("""
<style>
    /* Main Header Styling */
    .main-header {
        font-size: 3rem;
        color: #2563EB;
        text-align: center;
        margin-bottom: 2rem;
        font-weight: 700;
        letter-spacing: 1px;
        text-shadow: 2px 2px 4px rgba(0, 0, 0, 0.1);
        padding: 1rem;
        background: linear-gradient(90deg, #EFF6FF 0%, #DBEAFE 100%);
        border-radius: 12px;
        box-shadow: 0 4px 12px rgba(0, 0, 0, 0.1);
        transition: transform 0.3s ease;
    }
    .main-header:hover {
        transform: translateY(-5px);
    }

    /* Question Box Styling */
    .question-box {
        background-color: #F0F9FF;
        border-left: 5px solid #0EA5E9;
        padding: 20px;
        margin-bottom: 25px;
        border-radius: 8px;
        box-shadow: 0 4px 12px rgba(0, 0, 0, 0.08);
        transition: transform 0.2s ease, box-shadow 0.3s ease;
    }
    .question-box:hover {
        transform: translateY(-3px);
        box-shadow: 0 6px 16px rgba(0, 0, 0, 0.1);
    }

    /* Button Styling */
    .stButton>button {
        background-color: #2563EB;
        color: white;
        font-weight: bold;
        padding: 0.6rem 1.2rem;
        border-radius: 10px;
        transition: all 0.3s;
        border: none;
    }
    .stButton>button:hover {
        background-color: #1D4ED8;
        transform: translateY(-2px);
        box-shadow: 0 4px 12px rgba(0, 0, 0, 0.15);
    }

    /* Error Message Styling */
    .error-message {
        color: #DC2626;
        font-weight: bold;
        padding: 15px;
        background-color: #FEE2E2;
        border-radius: 8px;
        border-left: 5px solid #DC2626;
        margin: 10px 0;
    }

    /* Subheader Styling */
    .subheader {
        font-size: 1.8rem;
        color: #1E40AF;
        margin-top: 1.5rem;
        margin-bottom: 1rem;
        border-bottom: 3px solid #BFDBFE;
        padding-bottom: 0.6rem;
        font-weight: 600;
    }

    /* Question Number Styling */
    .question-number {
        font-weight: bold;
        color: #2563EB;
        font-size: 1.25rem;
    }

    /* Question Text Styling */
    .question-text {
        font-size: 1.15rem;
        margin-bottom: 1rem;
        line-height: 1.5;
    }

    /* Answer Styling */
    .answer {
        margin-top: 0.8rem;
        color: #047857;
        font-weight: bold;
        background-color: #ECFDF5;
        padding: 8px 12px;
        border-radius: 6px;
        display: inline-block;
        transition: background-color 0.3s ease;
    }
    .answer:hover {
        background-color: #D1FAE5;
    }

    /* Explanation Styling */
    .explanation {
        background-color: #F8FAFC;
        padding: 15px;
        border-radius: 8px;
        margin-top: 0.8rem;
        margin-bottom: 1.5rem;
        font-style: italic;
        border: 1px solid #E2E8F0;
        transition: background-color 0.3s ease;
    }
    .explanation:hover {
        background-color: #E2E8F0;
    }

    /* Stats Card Styling */
    .stats-card {
        background: linear-gradient(145deg, #DBEAFE 0%, #E0F2FE 100%);
        padding: 20px;
        border-radius: 12px;
        margin-bottom: 20px;
        box-shadow: 0 4px 6px rgba(0, 0, 0, 0.05);
        transition: transform 0.2s, box-shadow 0.3s ease;
    }
    .stats-card:hover {
        transform: translateY(-5px);
        box-shadow: 0 6px 12px rgba(0, 0, 0, 0.1);
    }

    .stats-card h4 {
        color: #1E40AF;
        font-size: 1.1rem;
        margin-bottom: 8px;
    }
    .stats-card h2 {
        color: #2563EB;
        font-size: 2.5rem;
        margin: 0;
    }

    /* File Uploader Styling */
    .file-uploader {
        background-color: #F1F5F9;
        padding: 20px;
        border-radius: 10px;
        border: 2px dashed #CBD5E1;
        transition: background-color 0.3s ease;
    }
    .file-uploader:hover {
        background-color: #E5E7EB;
    }

    /* Config Section Styling */
    .config-section {
        background-color: #F8FAFC;
        padding: 20px;
        border-radius: 10px;
        margin-bottom: 20px;
        box-shadow: 0 2px 8px rgba(0, 0, 0, 0.05);
        transition: background-color 0.3s ease;
    }
    .config-section:hover {
        background-color: #E2E8F0;
    }

    /* Action Button Styling */
    .action-button {
        margin-top: 8px;
    }

    /* Footer Styling */
    .footer {
        text-align: center;
        color: #64748B;
        font-size: 0.9rem;
        padding: 15px;
        margin-top: 30px;
        border-top: 1px solid #E2E8F0;
    }

    /* Custom Tab Styling */
    .stTabs [data-baseweb="tab-list"] {
        gap: 8px;
    }
    .stTabs [data-baseweb="tab"] {
        height: 50px;
        white-space: pre-wrap;
        background-color: #EFF6FF;
        border-radius: 8px 8px 0 0;
        gap: 1px;
        padding-top: 10px;
        padding-bottom: 10px;
    }
    .stTabs [aria-selected="true"] {
        background-color: #DBEAFE;
        border-bottom: 3px solid #3B82F6;
    }

    /* File List Styling */
    .file-list {
        max-height: 200px;
        overflow-y: auto;
        margin-top: 10px;
        padding: 10px;
        background-color: #F8FAFC;
        border-radius: 8px;
    }
    .file-item {
        display: flex;
        justify-content: space-between;
        padding: 8px;
        margin-bottom: 5px;
        background-color: white;
        border-radius: 6px;
        box-shadow: 0 1px 3px rgba(0, 0, 0, 0.1);
        transition: background-color 0.3s ease, box-shadow 0.2s ease;
    }
    .file-item:hover {
        background-color: #F3F4F6;
        box-shadow: 0 4px 8px rgba(0, 0, 0, 0.1);
    }

    .file-name {
        flex-grow: 1;
    }
    .file-size {
        color: #64748B;
        font-size: 0.8rem;
    }
    .remove-btn {
        color: #DC2626;
        cursor: pointer;
        margin-left: 10px;
    }
</style>
""", unsafe_allow_html=True)

# History tracking
if 'generation_history' not in st.session_state:
    st.session_state.generation_history = []

# File management
if 'uploaded_files' not in st.session_state:
    st.session_state.uploaded_files = []

def extract_text(file):
    try:
        # Read file content once
        file_bytes = file.getvalue()

        if file.type == "text/plain":
            return file_bytes.decode("utf-8")

        elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            from docx import Document
            doc = Document(BytesIO(file_bytes))
            return "\n".join([para.text for para in doc.paragraphs])

        elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            from pptx import Presentation
            prs = Presentation(BytesIO(file_bytes))
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)

        elif file.type == "application/pdf":
            import PyPDF2
            pdf_reader = PyPDF2.PdfReader(BytesIO(file_bytes))
            text = ""
            for page in pdf_reader.pages:
                if page.extract_text():
                    text += page.extract_text() + "\n"
            return text

        else:
            st.error("❌ Unsupported file type")
            return ""

    except Exception as e:
        st.error(f"Error processing file: {str(e)}")
        return ""


def format_questions(raw_questions):
    """Format the questions with better spacing and styling"""
    if not raw_questions:
        return ""
    
    # Split by numbered questions
    import re
    questions = re.split(r'\n\s*(\d+)\.\s+', raw_questions)
    
    if len(questions) <= 1:
        return raw_questions
    
    formatted = ""
    
    # Skip the first empty element if it exists
    start_idx = 1 if questions[0].strip() == "" else 0
    
    for i in range(start_idx, len(questions), 2):
        if i+1 < len(questions):
            question_num = questions[i]
            question_content = questions[i+1].strip()
            
            # Add extra spacing between parts
            question_content = question_content.replace("Question:", "\nQuestion:")
            question_content = question_content.replace("Options:", "\nOptions:")
            question_content = question_content.replace("Correct Answer:", "\n\nCorrect Answer:")
            question_content = question_content.replace("Explanation:", "\n\nExplanation:")
            
            formatted += f'<div class="question-box">\n'
            formatted += f'<span class="question-number">🔍 {question_num}.</span>\n'
            formatted += f'{question_content}\n'
            formatted += f'</div>\n\n'
    
    return formatted

def generate_with_gemini(text, question_type, difficulty, api_key, num_questions=5):
    try:
        genai.configure(api_key=api_key)
        
        # Use the correct model name - gemini-1.5-pro-latest or gemini-1.0-pro
        model = genai.GenerativeModel('gemini-1.5-flash')
        
        prompt = f"""
        Generate {num_questions} {difficulty.lower()} {question_type} questions based on this content.
        Format should be:
        
        1. Question: [question text]
           Options (if MCQ): A) [option1] B) [option2] C) [option3] D) [option4]
           Correct Answer: [correct answer]
           Explanation: [brief explanation]
        
        Content:
        {text[:12000]}  # Using first 12k chars to avoid token limits
        """
        
        response = model.generate_content(prompt)
        
        if not response.text:
            raise ValueError("Empty response from API")
            
        return response.text
    
    except Exception as e:
        st.markdown(f'<div class="error-message">⚠️ API error: {str(e)}</div>', unsafe_allow_html=True)
        st.info("🔑 Please ensure you're using the correct API key and model name")
        st.info("💡 Try using 'gemini-1.0-pro' if this model isn't available")
        return None

def remove_file(index):
    """Remove a file from the uploaded files list"""
    st.session_state.uploaded_files.pop(index)
    st.experimental_rerun()

def main():
    st.markdown('<h1 class="main-header">🧠 Lecture2Exam ✨</h1>', unsafe_allow_html=True)
    
    # Create tabs for main functionality and history
    tab1, tab2 = st.tabs(["📝 Generate Questions", "📊 History & Analytics"])
    
    with tab1:
        col1, col2 = st.columns([1, 2])
        
        with col1:
            st.markdown('<div class="config-section">', unsafe_allow_html=True)
            st.markdown("### 🔑 API Configuration")
            
            api_key = st.text_input("API Key", type="password", 
                                   help="Get your key from Google AI Studio")
            st.markdown('</div>', unsafe_allow_html=True)
            
            st.markdown('<div class="file-uploader">', unsafe_allow_html=True)
            st.markdown("### 📄 Upload Learning Material")
            
            # Multi-file uploader
            uploaded_files = st.file_uploader("Choose files", 
                                           type=["txt", "docx", "pptx", "pdf"],
                                           accept_multiple_files=True)
            
            # Store uploaded files in session state
            if uploaded_files and len(uploaded_files) > 0:
                for file in uploaded_files:
                    if file.name not in [f['name'] for f in st.session_state.uploaded_files]:
                        st.session_state.uploaded_files.append({
                            'name': file.name,
                            'size': f"{len(file.getvalue()) / 1024:.1f} KB",
                            'file_obj': file
                        })
            
            # Display uploaded files with remove option
            if len(st.session_state.uploaded_files) > 0:
                st.markdown("### 📂 File History")
                st.markdown('<div class="file-list">', unsafe_allow_html=True)
                for i, file in enumerate(st.session_state.uploaded_files):
                    st.markdown(
                        f'<div class="file-item">'
                        f'<span class="file-name">{file["name"]}</span>'
                        f'<span class="file-size">{file["size"]}</span>'
                        f'<span class="remove-btn" onclick="removeFile({i})"></span>'
                        f'</div>',
                        unsafe_allow_html=True
                    )
                st.markdown('</div>', unsafe_allow_html=True)
            
            st.markdown('</div>', unsafe_allow_html=True)
            
            if len(st.session_state.uploaded_files) > 0:
                st.markdown('<div class="config-section">', unsafe_allow_html=True)
                st.markdown("### ⚙️ Question Settings")
                
                question_type = st.selectbox("Question Type 📋", 
                                           ["MCQ", "Fill-in-the-Blank", "Short Answer"])
                
                difficulty_icons = {"Easy": "🟢", "Medium": "🟡", "Hard": "🔴"}
                difficulties = [f"{v} {k}" for k, v in difficulty_icons.items()]
                difficulty_display = st.selectbox("Difficulty Level", difficulties)
                difficulty = difficulty_display.split(' ')[1]  # Extract the actual difficulty
                
                num_questions = st.slider("Number of Questions 🔢", 
                                        min_value=3, max_value=15, value=5,
                                        help="More questions will take longer to generate")
                st.markdown('</div>', unsafe_allow_html=True)
                
                st.markdown('<div class="action-button">', unsafe_allow_html=True)
                gen_col1, gen_col2 = st.columns(2)
                with gen_col1:
                    generate_btn = st.button("✨ Generate Questions", use_container_width=True)
                with gen_col2:
                    clear_btn = st.button("🗑️ Clear Results", use_container_width=True)
                st.markdown('</div>', unsafe_allow_html=True)
                
                if clear_btn:
                    if 'questions' in st.session_state:
                        del st.session_state.questions
                
                if generate_btn:
                    if not api_key:
                        st.error("⚠️ Please enter your API key")
                    else:
                        with st.spinner("🧙‍♂️ Generating intelligent questions..."):
                            # Combine all files' content
                            combined_text = ""
                            for file_info in st.session_state.uploaded_files:
                                file = file_info['file_obj']
                                text = extract_text(file)
                                if text:
                                    combined_text += f"\n\n--- {file.name} ---\n\n{text}"
                            
                            if combined_text:
                                # Try with latest model first, fallback to 1.0 if needed
                                questions = generate_with_gemini(combined_text, question_type, difficulty, api_key, num_questions)
                                if not questions:
                                    # Fallback to gemini-1.0-pro if latest model fails
                                    genai.configure(api_key=api_key)
                                    model = genai.GenerativeModel('gemini-1.0-pro')
                                    questions = generate_with_gemini(combined_text, question_type, difficulty, api_key, num_questions)
                                
                                if questions:
                                    st.session_state.questions = questions
                                    st.session_state.formatted_questions = format_questions(questions)
                                    
                                    # Add to history
                                    file_names = ", ".join([f['name'] for f in st.session_state.uploaded_files])
                                    st.session_state.generation_history.append({
                                        'timestamp': datetime.now().strftime("%Y-%m-%d %H:%M"),
                                        'file': file_names,
                                        'type': question_type,
                                        'difficulty': difficulty,
                                        'count': num_questions,
                                        'questions': questions
                                    })
                                else:
                                    st.session_state.questions = None
        
        with col2:
            if len(st.session_state.uploaded_files) > 0 and 'questions' in st.session_state and st.session_state.questions:
                st.markdown(f"<h2 class='subheader'>✅ Generated {question_type} Questions</h2>", unsafe_allow_html=True)
                
                # Show formatted questions with improved spacing
                if 'formatted_questions' in st.session_state:
                    st.markdown(st.session_state.formatted_questions, unsafe_allow_html=True)
                else:
                    st.markdown(st.session_state.questions)
                
                col_a, col_b = st.columns(2)
                with col_a:
                    st.download_button(
                        label="📥 Download Questions",
                        data=BytesIO(st.session_state.questions.encode('utf-8')),
                        file_name=f"{question_type.lower().replace(' ', '_')}_questions.txt",
                        mime="text/plain",
                        use_container_width=True
                    )
                with col_b:
                    if st.button("🔄 Regenerate Questions", use_container_width=True):
                        with st.spinner("🔄 Regenerating questions..."):
                            # Combine all files' content again
                            combined_text = ""
                            for file_info in st.session_state.uploaded_files:
                                file = file_info['file_obj']
                                text = extract_text(file)
                                if text:
                                    combined_text += f"\n\n--- {file.name} ---\n\n{text}"
                            
                            if combined_text:
                                questions = generate_with_gemini(combined_text, question_type, difficulty, api_key, num_questions)
                                if questions:
                                    st.session_state.questions = questions
                                    st.session_state.formatted_questions = format_questions(questions)
            else:
                st.markdown("""
                <div style="background-color: #F0F9FF; padding: 30px; border-radius: 10px; text-align: center; margin-top: 50px;">
                    <h3>🚀 Ready to Create Questions?</h3>
                    <p>Upload your learning material and configure settings to generate AI-powered questions.</p>
                    <p style="font-size: 3rem; margin: 20px 0;">📚 ➡️ 🧠 ➡️ 📝</p>
                    <p>Perfect for teachers, students, and learning professionals.</p>
                </div>
                """, unsafe_allow_html=True)
    
    with tab2:
        st.markdown("<h2 class='subheader'>📊 Generation History & Analytics</h2>", unsafe_allow_html=True)
        
        if not st.session_state.generation_history:
            st.info("📭 No question generation history yet. Generate some questions to see your history.")
        else:
            history_df = pd.DataFrame(st.session_state.generation_history)
            
            # Show stats
            st.markdown("<h3 class='subheader'>📈 Usage Statistics</h3>", unsafe_allow_html=True)
            col1, col2, col3 = st.columns(3)
            with col1:
                st.markdown(f"""
                <div class="stats-card">
                <h4>🔄 Total Generations</h4>
                <h2>{len(history_df)}</h2>
                </div>
                """, unsafe_allow_html=True)
            with col2:
                st.markdown(f"""
                <div class="stats-card">
                <h4>📋 Question Types</h4>
                <h2>{len(history_df['type'].unique())}</h2>
                </div>
                """, unsafe_allow_html=True)
            with col3:
                st.markdown(f"""
                <div class="stats-card">
                <h4>📑 Files Processed</h4>
                <h2>{len(history_df['file'].unique())}</h2>
                </div>
                """, unsafe_allow_html=True)
            
            # Additional stats row (modified)
            if len(history_df) > 0:
                col4 = st.columns(1)[0]
                with col4:
                    total_questions = history_df['count'].sum()
                    st.markdown(f"""
                    <div class="stats-card">
                    <h4>🧩 Total Questions Created</h4>
                    <h2>{total_questions}</h2>
                    </div>
                    """, unsafe_allow_html=True)
                
                # New analytics section
                st.markdown("<h3 class='subheader'>📊 Detailed Analytics</h3>", unsafe_allow_html=True)
                
                # Prepare data for visualizations
                history_df['datetime'] = pd.to_datetime(history_df['timestamp'])
                timeline_df = history_df.set_index('datetime').resample('D')['count'].sum().reset_index()
                difficulty_df = history_df.groupby('difficulty')['count'].sum().reset_index()
                
                # Create two columns for charts
                chart_col1, chart_col2 = st.columns(2)
                
                with chart_col1:
                    # Timeline chart
                    fig = px.line(timeline_df, 
                                x='datetime', 
                                y='count',
                                title='Question Generation Timeline',
                                labels={'count': 'Questions Generated', 'datetime': 'Date'},
                                markers=True)
                    fig.update_layout(hovermode="x unified")
                    st.plotly_chart(fig, use_container_width=True)
                
                with chart_col2:
                    # Difficulty distribution pie chart
                    fig = px.pie(difficulty_df,
                                values='count',
                                names='difficulty',
                                title='Difficulty Distribution',
                                color_discrete_sequence=px.colors.qualitative.Pastel,
                                hole=0.3)
                    fig.update_traces(textposition='inside', textinfo='percent+label')
                    st.plotly_chart(fig, use_container_width=True)
            
            # Display history table
            st.markdown("<h3 class='subheader'>📜 Recent Generation Activity</h3>", unsafe_allow_html=True)
            
            # Create a more readable history table
            display_df = history_df[['timestamp', 'file', 'type', 'difficulty', 'count']].copy()
            display_df.columns = ['⏰ Timestamp', '📄 File', '📋 Type', '🎯 Difficulty', '🔢 Count']
            st.dataframe(display_df.tail(10), use_container_width=True)
            
            # Option to clear history
            if st.button("🗑️ Clear History"):
                st.session_state.generation_history = []
                st.experimental_rerun()

    # Footer
    st.markdown("---")
    st.markdown(
        """
        <div class="footer">
            🧠 Lecture2Exam<br>
            Make learning more effective with AI-powered assessments ✨<br>
            Made with teamwork of Shreyas, Shaurya and Mahati 🎯
        </div>
        """, 
        unsafe_allow_html=True
    )
    
    # Add JavaScript for file removal
    st.markdown("""
    <script>
    function removeFile(index) {
        window.parent.postMessage({
            type: 'streamlit:setComponentValue',
            value: index
        }, '*');
    }
    </script>
    """, unsafe_allow_html=True)

if __name__ == "__main__":
    main()
