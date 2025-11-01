import streamlit as st
from io import BytesIO
import google.generativeai as genai
import pandas as pd
from datetime import datetime
import plotly.express as px
import re

# Set page configuration
st.set_page_config(
    page_title="Lecture2Exam",
    page_icon="🧠",
    layout="wide",
    initial_sidebar_state="collapsed"
)

st.markdown("""
<style>
    /* Enhanced Main Header with Animation */
    .main-header {
        font-size: 3.2rem;
        color: #2563EB;
        text-align: center;
        margin-bottom: 2rem;
        font-weight: 700;
        background: linear-gradient(120deg, #3B82F6 0%, #1E40AF 100%);
        background-clip: text;
        -webkit-background-clip: text;
        color: transparent;
        padding: 1.2rem 0;
        position: relative;
        overflow: hidden;
    }
    .main-header::after {
        content: "";
        position: absolute;
        bottom: 0;
        left: 0;
        width: 100%;
        height: 4px;
        background: linear-gradient(90deg, #3B82F6, #60A5FA, #3B82F6);
        background-size: 200% 100%;
        animation: gradient-shift 3s ease infinite;
    }
    @keyframes gradient-shift {
        0% {background-position: 0% 50%}
        50% {background-position: 100% 50%}
        100% {background-position: 0% 50%}
    }
    
    /* Interactive File Cards */
    .file-card {
        background: white;
        border-radius: 10px;
        box-shadow: 0 4px 12px rgba(0,0,0,0.08);
        padding: 15px;
        margin-bottom: 15px;
        display: flex;
        justify-content: space-between;
        align-items: center;
        transition: all 0.3s ease;
        border-left: 4px solid #3B82F6;
    }
    .file-card:hover {
        transform: translateY(-5px);
        box-shadow: 0 8px 15px rgba(0,0,0,0.12);
    }
    
    /* Improved Question Box with Animation */
    .question-box {
        background: linear-gradient(to right, #F0F9FF 0%, #EFF6FF 100%);
        border-left: 5px solid #0EA5E9;
        padding: 20px;
        margin-bottom: 25px;
        border-radius: 8px;
        box-shadow: 0 4px 12px rgba(0,0,0,0.08);
        transition: all 0.3s cubic-bezier(0.25, 0.8, 0.25, 1);
        position: relative;
        overflow: hidden;
    }
    .question-box:hover {
        transform: translateY(-3px) scale(1.01);
        box-shadow: 0 8px 20px rgba(0,0,0,0.12);
    }
    .question-box::before {
        content: "";
        position: absolute;
        top: 0;
        left: 0;
        width: 100%;
        height: 100%;
        background: linear-gradient(120deg, rgba(59, 130, 246, 0.1) 0%, rgba(14, 165, 233, 0.05) 100%);
        transform: translateX(-100%);
        transition: transform 0.5s ease;
    }
    .question-box:hover::before {
        transform: translateX(0);
    }
    
    /* Pulsing Action Button */
    .action-button button {
        position: relative;
        overflow: hidden;
    }
    .action-button button::after {
        content: "";
        position: absolute;
        top: 50%;
        left: 50%;
        width: 5px;
        height: 5px;
        background: rgba(255, 255, 255, 0.5);
        opacity: 0;
        border-radius: 100%;
        transform: scale(1, 1) translate(-50%, -50%);
        transform-origin: 50% 50%;
    }
    .action-button button:hover::after {
        animation: ripple 1s ease-out;
    }
    @keyframes ripple {
        0% {
            transform: scale(0, 0);
            opacity: 0.5;
        }
        100% {
            transform: scale(20, 20);
            opacity: 0;
        }
    }
</style>
""", unsafe_allow_html=True)

def extract_text(file):
    """Extract text from various file formats"""
    try:
        if file.type == "text/plain":
            file.seek(0)
            return file.getvalue().decode("utf-8")
        elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            from docx import Document
            file.seek(0)
            doc = Document(BytesIO(file.read()))
            return "\n".join([para.text for para in doc.paragraphs])
        elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            from pptx import Presentation
            file.seek(0)
            prs = Presentation(BytesIO(file.read()))
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)
        elif file.type == "application/pdf":
            import PyPDF2
            file.seek(0)
            pdf_reader = PyPDF2.PdfReader(BytesIO(file.read()))
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
            return text
        else:
            st.error(f"Unsupported file type: {file.type}")
            return ""
    except Exception as e:
        st.error(f"Error processing file: {str(e)}")
        return ""

def format_questions(raw_questions):
    """Format the questions with better spacing and styling"""
    if not raw_questions:
        return ""

    # Split by numbered questions (updated regex)
    questions = re.split(r'(?:\n|^)\s*(\d+)\.\s+', raw_questions)

    if len(questions) <= 1:
        return raw_questions

    formatted = ""
    start_idx = 1 if questions[0].strip() == "" else 0

    for i in range(start_idx, len(questions), 2):
        if i+1 < len(questions):
            question_num = questions[i]
            question_content = questions[i+1].strip()

            # Enhanced formatting logic
            question_content = re.sub(
                r'(Question:|Options:|Correct Answer:|Explanation:)',
                r'<strong>\1</strong>', 
                question_content
            )
            question_content = re.sub(
                r'([A-D])\) (.*?)(?=\s*[A-D]\)|$|Correct Answer:)',
                r'<span style="margin-left:15px;display:block;"><span style="color:#3B82F6;font-weight:bold;">\1)</span> \2</span>',
                question_content
            )
            question_content = re.sub(
                r'<strong>Correct Answer:</strong>\s*([A-D])',
                r'<strong>Correct Answer:</strong> <span style="color:#059669;font-weight:bold;">\1</span>',
                question_content
            )

            formatted += f'<div class="question-box">\n'
            formatted += f'<span class="question-number">🔍 {question_num}.</span>\n'
            formatted += f'{question_content}\n'
            formatted += f'</div>\n\n'

    return formatted

# Enhanced file management
def handle_multiple_files():
    # File uploading interface
    st.markdown('<div class="file-uploader">', unsafe_allow_html=True)
    st.markdown("### 📄 Upload Learning Materials")

    uploaded_files = st.file_uploader("Choose files", 
                                   type=["txt", "docx", "pptx", "pdf"],
                                   accept_multiple_files=True)

    # Initialize session state for file storage if it doesn't exist
    if 'uploaded_files' not in st.session_state:
        st.session_state.uploaded_files = []

    # Process newly uploaded files
    if uploaded_files:
        # Clear existing files list first to avoid duplicates
        seen_files = {f['name']: f for f in st.session_state.uploaded_files}

        for file in uploaded_files:
            # Reset file position to beginning to ensure proper reading
            file.seek(0)

            # Extract text and check if it's empty
            preview_text = extract_text(file)
            if not preview_text or preview_text.strip() == "":
                st.warning(f"No content could be extracted from {file.name}. Please check the file.")
                continue

            # Generate preview for display
            preview = preview_text[:100] + "..." if len(preview_text) > 100 else preview_text

            # Store the file with its content
            seen_files[file.name] = {
                'name': file.name,
                'size': f"{len(file.getvalue()) / 1024:.1f} KB",
                'hash': hash(file.name + str(file.size)),
                'file_obj': file,
                'preview': preview,
                'content': preview_text,  # Store the full content
                'timestamp': datetime.now().strftime("%Y-%m-%d %H:%M"),
                'type': file.type
            }

        st.session_state.uploaded_files = list(seen_files.values())

    # Display uploaded files with interactive cards
    if st.session_state.uploaded_files:
        st.markdown("### 📂 Your Materials")

        files_to_remove = []  # Track files to remove

        for i, file in enumerate(st.session_state.uploaded_files):
            col1, col2 = st.columns([4, 1])

            with col1:
                expander = st.expander(f"📄 {file['name']} ({file['size']})")
                with expander:
                    st.write(f"**Type:** {file['type'].split('/')[-1].upper()}")
                    st.write(f"**Added:** {file['timestamp']}")
                    st.write("**Preview:**")
                    st.markdown(f"<div style='background:#f5f5f5;padding:10px;border-radius:5px;font-family:monospace;white-space:pre-wrap;'>{file['preview']}</div>", unsafe_allow_html=True)

                    # Show content length info
                    content_len = len(file.get('content', ''))
                    if content_len > 0:
                        st.success(f"✅ Extracted {content_len} characters of text")
                    else:
                        st.error("❌ No content could be extracted")

            with col2:
                # Fixed: Individual file removal
                if st.button(f"🗑️ Remove", key=f"remove_{file['name']}_{i}"):
                    files_to_remove.append(i)

        # Remove files that were marked for deletion
        if files_to_remove:
            for index in sorted(files_to_remove, reverse=True):
                st.session_state.uploaded_files.pop(index)
            st.rerun()
            st.rerun()

    st.markdown('</div>', unsafe_allow_html=True)

    # Return true if files are available
    return len(st.session_state.uploaded_files) > 0

def generate_questions(api_key, question_type, difficulty, num_questions):
    with st.spinner("🧙‍♂️ Creating intelligent questions..."):
        # Combine content from all files with better organization
        combined_text = ""

        # Check if there are any files with content
        has_content = False
        for file_info in st.session_state.uploaded_files:
            # Use the stored content instead of re-extracting
            text = file_info.get('content', '')
            if text and text.strip():
                has_content = True
                # Add better file context separation
                combined_text += f"\n\n==== CONTENT FROM: {file_info['name']} ====\n\n{text}\n\n"

        if not has_content:
            st.error("⚠️ Could not extract text from any of the files. Please check your files.")
            return None

        # Configure Gemini AI
        genai.configure(api_key=api_key)

        try:
            return genai.GenerativeModel('gemini-1.5-flash')
        except:
            return genai.GenerativeModel('gemini-1.0-pro')

            # Enhanced prompt for better question quality
            prompt = f"""
            Generate {num_questions} {difficulty.lower()} {question_type} questions based on the content I'm providing.
            
            Guidelines:
            - Focus on key concepts and important details
            - Ensure questions test understanding, not just memorization
            - Include a mix of recall and application questions
            - Make explanations clear and educational
            
            Format each question as:
            
            [#]. Question: [question text]
                Options (for MCQ): A) [option1] B) [option2] C) [option3] D) [option4]
                Correct Answer: [correct answer]
                Explanation: [concise explanation]
            
            Content:
            {combined_text[:15000]}
            """

            response = model.generate_content(prompt)

            if not response.text:
                # Fallback to gemini-1.0-pro if needed
                model = genai.GenerativeModel('gemini-1.0-pro')
                response = model.generate_content(prompt)

            return response.text

        except Exception as e:
            st.markdown(f'<div class="error-message">⚠️ API error: {str(e)}</div>', unsafe_allow_html=True)
            return None

def display_history_analytics():
    """Display history and analytics with proper data handling"""
    st.markdown("<h2 class='subheader'>📊 Generation History & Analytics</h2>", unsafe_allow_html=True)

    if 'generation_history' not in st.session_state or not st.session_state.generation_history:
        st.info("📭 No question generation history yet. Generate some questions to see your history.")
        return

    history_df = pd.DataFrame(st.session_state.generation_history)

    # Show stats
    st.markdown("<h3 class='subheader'>📈 Usage Statistics</h3>", unsafe_allow_html=True)
    col1, col2, col3 = st.columns(3)
    with col1:
        st.markdown(f"""
        <div class="stats-card">
        <h4>🔄 Total Generations</h4>
        <h2>{len(history_df)}</h2>
        </div>
        """, unsafe_allow_html=True)
    with col2:
        st.markdown(f"""
        <div class="stats-card">
        <h4>📋 Question Types</h4>
        <h2>{history_df['type'].nunique()}</h2>
        </div>
        """, unsafe_allow_html=True)
    with col3:
        st.markdown(f"""
        <div class="stats-card">
        <h4>📑 Files Processed</h4>
        <h2>{history_df['file'].nunique()}</h2>
        </div>
        """, unsafe_allow_html=True)

    # Additional stats
    total_questions = history_df['count'].sum()
    st.markdown(f"""
    <div class="stats-card" style="margin-top:15px;">
    <h4>🧩 Total Questions Created</h4>
    <h2>{total_questions}</h2>
    </div>
    """, unsafe_allow_html=True)

    # Analytics visualizations
    st.markdown("<h3 class='subheader'>📊 Detailed Analytics</h3>", unsafe_allow_html=True)
    
    try:
        # Prepare data for visualizations
        history_df['datetime'] = pd.to_datetime(history_df['timestamp'])
        timeline_df = history_df.set_index('datetime').resample('D')['count'].sum().reset_index()
        difficulty_df = history_df.groupby('difficulty')['count'].sum().reset_index()

        # Create charts
        chart_col1, chart_col2 = st.columns(2)
        with chart_col1:
            fig = px.line(timeline_df, 
                        x='datetime', 
                        y='count',
                        title='Question Generation Timeline',
                        labels={'count': 'Questions Generated', 'datetime': 'Date'},
                        markers=True)
            fig.update_layout(hovermode="x unified")
            st.plotly_chart(fig, use_container_width=True)

        with chart_col2:
            fig = px.pie(difficulty_df,
                        values='count',
                        names='difficulty',
                        title='Difficulty Distribution',
                        color_discrete_sequence=px.colors.qualitative.Pastel,
                        hole=0.3)
            fig.update_traces(textposition='inside', textinfo='percent+label')
            st.plotly_chart(fig, use_container_width=True)

    except Exception as e:
        st.error(f"Error generating charts: {str(e)}")

    # Display history table
    st.markdown("<h3 class='subheader'>📜 Recent Generation Activity</h3>", unsafe_allow_html=True)
    display_df = history_df[['timestamp', 'file', 'type', 'difficulty', 'count']].copy()
    display_df.columns = ['⏰ Timestamp', '📄 File', '📋 Type', '🎯 Difficulty', '🔢 Count']
    st.dataframe(display_df.tail(10), use_container_width=True)

    if st.button("🗑️ Clear History", key="clear_history"):
        st.session_state.generation_history = []
        st.rerun()

def main():
    # Initialize session state
    if 'generation_history' not in st.session_state:
        st.session_state.generation_history = []
    if 'uploaded_files' not in st.session_state:
        st.session_state.uploaded_files = []
    if 'questions' not in st.session_state:
        st.session_state.questions = None
    if 'formatted_questions' not in st.session_state:
        st.session_state.formatted_questions = ""

    # App header
    st.markdown('<h1 class="main-header">🧠 Lecture2Exam ✨</h1>', unsafe_allow_html=True)

    # Create tabs
    tab1, tab2 = st.tabs(["📝 Create Questions", "📊 History & Analytics"])

    with tab1:
        col1, col2 = st.columns([1, 2])

        with col1:
            # API Configuration
            with st.container():
                st.subheader("🔑 API Configuration")
    
                api_key = st.text_input(
                    label="API Key",
                    type="password",
                    placeholder="Enter your API key here...",
                    help="This key is required to access the AI services. Keep it confidential."
                  )                    
                st.markdown("You may obtain API Key from [Google AI Studio](https://aistudio.google.com/apikey).")
            # File handling
            files_available = handle_multiple_files()

            # Question settings
            if files_available:
                with st.container():
                    st.markdown("### ⚙️ Question Settings")
                    question_type = st.selectbox("Question Type", 
                                               ["MCQ", "Fill-in-the-Blank", "Short Answer"])
                    difficulty = st.select_slider(
                        "Difficulty Level",
                        options=["Easy", "Medium", "Hard"],
                        format_func=lambda x: f"{'🟢' if x == 'Easy' else '🟡' if x == 'Medium' else '🔴'} {x}"
                    )
                    num_questions = st.slider("Number of Questions", 
                                            min_value=3, max_value=15, value=5)

                # Generate button
                st.markdown('<div class="action-button">', unsafe_allow_html=True)
                generate_btn = st.button("✨ Generate Questions", use_container_width=True)
                st.markdown('</div>', unsafe_allow_html=True)

                if generate_btn:
                    if not api_key:
                        st.error("⚠️ Please enter your API key")
                    else:
                        questions = generate_questions(api_key, question_type, difficulty, num_questions)
                        if questions:
                            file_names = ", ".join([f['name'] for f in st.session_state.uploaded_files])
                            st.session_state.questions = questions
                            st.session_state.formatted_questions = format_questions(questions)
                            history_entry = {
                    'timestamp': datetime.now().isoformat(),  # ISO format for pandas parsing
                    'file' : file_names,
                    'type': question_type,
                    'difficulty': difficulty,
                    'count': num_questions,
                    'raw_content': questions  # Store original generated text
                }
                
                # Update history (keep last 50 entries)
                            st.session_state.generation_history = [history_entry] + st.session_state.generation_history[:49]



        with col2:
            if files_available and st.session_state.questions:
                st.markdown(f"<h2 class='subheader'>✅ Generated {question_type} Questions</h2>", unsafe_allow_html=True)
                st.markdown(st.session_state.formatted_questions, unsafe_allow_html=True)
                col_a, col_b = st.columns(2)
                with col_a:
                    st.download_button(
                        label="📥 Download Questions",
                        data=BytesIO(st.session_state.questions.encode('utf-8')),
                        file_name=f"{question_type.lower().replace(' ', '_')}_questions.txt",
                        mime="text/plain",
                        use_container_width=True
                    )
                with col_b:
                    if st.button("🔄 Regenerate Questions", use_container_width=True):
                        with st.spinner("🔄 Regenerating questions..."):
                            # Combine all files' content again
                            combined_text = ""
                            for file_info in st.session_state.uploaded_files:
                                file = file_info['file_obj']
                                text = extract_text(file)
                                if text:
                                    combined_text += f"\n\n--- {file.name} ---\n\n{text}"
                            
                            if combined_text:
                                questions = generate_questions(api_key, question_type, difficulty, num_questions)
                                if questions:
                                    file_names = ", ".join([f['name'] for f in st.session_state.uploaded_files])
                                    st.session_state.questions = questions
                                    st.session_state.formatted_questions = format_questions(questions)
                                    file_names = ", ".join([f['name'] for f in st.session_state.uploaded_files])
                                    st.session_state.generation_history.append({
                            'timestamp': datetime.now().isoformat(),
                            'file': file_names,
                            'type': question_type,
                            'difficulty': difficulty,
                            'count': num_questions,
                            'raw_questions': questions
                        })          
                                    st.rerun()
            else:
                st.markdown("""
                <div style="background-color: #F0F9FF; padding: 30px; border-radius: 10px; text-align: center; margin-top: 50px;">
                    <h3>🚀 Ready to Create Questions?</h3>
                    <p>Upload your learning material and configure settings to generate AI-powered questions.</p>
                    <p style="font-size: 3rem; margin: 20px 0;">📚 ➡️ 🧠 ➡️ 📝</p>
                    <p>Perfect for teachers, students, and learning professionals.</p>
                </div>
                """, unsafe_allow_html=True)

    

    with tab2:
        display_history_analytics()

    # Footer
    st.markdown("""
    <div class="footer" style="text-align: center; margin-top: 40px; padding: 20px; 
                             border-top: 1px solid #E2E8F0; animation: fadeIn 1s ease-in;">
        🧠 Lecture2Exam<br>
        Making assessment creation effortless with AI ✨<br>
        <small>Made with ❤️ by Shreyas</small>
    </div>
    """, unsafe_allow_html=True)

if __name__ == "__main__":
    main()
